import pytest
import tempfile
import os
from parsers.text import TextParser
from parsers.docx_parser import DocxParser
from parsers.xlsx_parser import XlsxParser
from parsers.pptx_parser import PptxParser


class TestTextParser:
    def test_parse_md(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as f:
            f.write("# Hello\n張三 test@test.com")
            path = f.name
        try:
            text = TextParser().parse(path)
            assert "張三" in text
            assert "test@test.com" in text
        finally:
            os.unlink(path)

    def test_parse_csv(self):
        with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False, encoding='utf-8') as f:
            f.write("name,email\n張三,test@test.com")
            path = f.name
        try:
            text = TextParser().parse(path)
            assert "張三" in text
        finally:
            os.unlink(path)


class TestDocxParser:
    def test_parse_docx(self):
        from docx import Document
        doc = Document()
        doc.add_paragraph("張三是教授")
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            path = f.name
        try:
            text = DocxParser().parse(path)
            assert "張三是教授" in text
        finally:
            os.unlink(path)


class TestXlsxParser:
    def test_parse_xlsx(self):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = '張三'
        ws['B1'] = 'test@test.com'
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            wb.save(f.name)
            path = f.name
        try:
            text = XlsxParser().parse(path)
            assert "張三" in text
            assert "test@test.com" in text
        finally:
            os.unlink(path)


class TestPptxParser:
    def test_parse_pptx(self):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "張三的報告"
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            path = f.name
        try:
            text = PptxParser().parse(path)
            assert "張三的報告" in text
        finally:
            os.unlink(path)
