from pptx import Presentation


class PptxParser:
    EXTENSIONS = {'.pptx'}
    OUTPUT_EXTENSION = ".pptx"

    def parse(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
        except Exception:
            return ""
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        parts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text)
        return '\n'.join(parts)

    def anonymize_to_path(self, file_path: str, output_path: str, anonymizer):
        try:
            prs = Presentation(file_path)
        except Exception:
            return False, []

        spans = []

        def iter_shapes(shapes):
            for shape in shapes:
                yield shape
                if hasattr(shape, "shapes"):
                    yield from iter_shapes(shape.shapes)

        for slide in prs.slides:
            for shape in iter_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    anonymized, shape_spans = anonymizer._anonymize_text_with_spans(shape.text)
                    if shape_spans:
                        shape.text = anonymized
                        spans.extend(shape_spans)

                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            anonymized, cell_spans = anonymizer._anonymize_text_with_spans(cell.text)
                            if cell_spans:
                                cell.text = anonymized
                                spans.extend(cell_spans)

        if not spans:
            return False, []

        prs.save(output_path)
        return True, spans
